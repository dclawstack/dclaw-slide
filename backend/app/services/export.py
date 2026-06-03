"""Export a Presentation to HTML zip, PPTX, or PDF.

All exports are pure-Python (no native deps) so they work the same on
local Mac dev, Linux CI, and the production python:3.11-slim image.
"""

from __future__ import annotations

import io
import re
import zipfile
from dataclasses import dataclass
from html import escape

# theme_accent is user-controlled and interpolated into a CSS string, so it must
# be restricted to a safe color token (e.g. "#0ea5e9") to prevent CSS/HTML
# injection. Anything that doesn't match falls back to this default.
_DEFAULT_ACCENT = "#0ea5e9"
_HEX_COLOR_RE = re.compile(r"#[0-9A-Fa-f]{3,8}")


def _safe_accent(theme_accent: str | None) -> str:
    candidate = (theme_accent or "").strip()
    if not _HEX_COLOR_RE.fullmatch(candidate):
        return _DEFAULT_ACCENT
    return candidate

from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import landscape, letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
)


@dataclass
class ExportSlide:
    position: int
    title: str
    body: str
    layout: str
    speaker_notes: str


@dataclass
class ExportInput:
    title: str
    theme_accent: str
    slides: list[ExportSlide]


# ──────────────────────────────────────────────────────────────────────────────
# HTML zip
# ──────────────────────────────────────────────────────────────────────────────

_HTML_TEMPLATE = """<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <title>{title}</title>
  <style>
    body {{ font-family: -apple-system, BlinkMacSystemFont, 'Inter', sans-serif;
            background: #0f172a; color: #f8fafc; margin: 0; padding: 0; }}
    .slide {{ width: 100vw; height: 100vh; display: flex; flex-direction: column;
              justify-content: center; padding: 4rem 6rem; box-sizing: border-box;
              page-break-after: always; border-bottom: 6px solid {accent}; }}
    h1 {{ font-size: 3.5rem; margin: 0 0 1.5rem; color: {accent}; }}
    .body {{ font-size: 1.4rem; white-space: pre-wrap; line-height: 1.6; }}
    .notes {{ margin-top: 2rem; padding-top: 1rem; border-top: 1px solid #334155;
              font-size: 0.85rem; color: #94a3b8; font-style: italic; }}
    .nav {{ position: fixed; top: 1rem; left: 1rem; font-size: 0.85rem;
            color: #94a3b8; }}
    .layout {{ position: fixed; top: 1rem; right: 1rem; font-size: 0.7rem;
               color: #475569; letter-spacing: 0.1em; text-transform: uppercase; }}
  </style>
</head>
<body>
{slides_html}
</body>
</html>
"""


def render_html(presentation: ExportInput) -> bytes:
    slides_html = []
    total = len(presentation.slides)
    for slide in presentation.slides:
        notes_html = (
            f'<div class="notes">{escape(slide.speaker_notes)}</div>'
            if slide.speaker_notes
            else ""
        )
        slides_html.append(
            f"""
  <section class="slide" id="slide-{slide.position + 1}">
    <div class="nav">{slide.position + 1} / {total}</div>
    <div class="layout">{escape(slide.layout)}</div>
    <h1>{escape(slide.title)}</h1>
    <div class="body">{escape(slide.body)}</div>
    {notes_html}
  </section>"""
        )
    document = _HTML_TEMPLATE.format(
        title=escape(presentation.title),
        accent=_safe_accent(presentation.theme_accent),
        slides_html="\n".join(slides_html),
    )

    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("index.html", document)
        zf.writestr("README.txt", "Open index.html in any browser.")
    return buffer.getvalue()


# ──────────────────────────────────────────────────────────────────────────────
# PPTX
# ──────────────────────────────────────────────────────────────────────────────


def render_pptx(presentation: ExportInput) -> bytes:
    pptx = PPTXPresentation()
    pptx.slide_width = Inches(13.333)
    pptx.slide_height = Inches(7.5)

    for slide_data in presentation.slides:
        # layout 6 = blank in default PPTX template
        slide = pptx.slides.add_slide(pptx.slide_layouts[6])

        # Title textbox
        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.2)
        )
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        para = title_tf.paragraphs[0]
        para.text = slide_data.title or " "
        for run in para.runs:
            run.font.size = Pt(40)
            run.font.bold = True

        # Body textbox
        body_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.0)
        )
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        first = True
        for line in (slide_data.body or "").splitlines():
            paragraph = body_tf.paragraphs[0] if first else body_tf.add_paragraph()
            paragraph.text = line
            for run in paragraph.runs:
                run.font.size = Pt(22)
            first = False
        if first:
            body_tf.paragraphs[0].text = ""

        # Speaker notes
        if slide_data.speaker_notes:
            slide.notes_slide.notes_text_frame.text = slide_data.speaker_notes

    buffer = io.BytesIO()
    pptx.save(buffer)
    return buffer.getvalue()


# ──────────────────────────────────────────────────────────────────────────────
# PDF
# ──────────────────────────────────────────────────────────────────────────────


def render_pdf(presentation: ExportInput) -> bytes:
    buffer = io.BytesIO()
    page_size = landscape(letter)
    doc = SimpleDocTemplate(
        buffer,
        pagesize=page_size,
        leftMargin=0.6 * inch,
        rightMargin=0.6 * inch,
        topMargin=0.5 * inch,
        bottomMargin=0.5 * inch,
        title=presentation.title,
    )
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "DeckTitle",
        parent=styles["Title"],
        fontSize=28,
        leading=34,
        spaceAfter=18,
        textColor="#0f172a",
    )
    body_style = ParagraphStyle(
        "DeckBody",
        parent=styles["BodyText"],
        fontSize=14,
        leading=20,
    )
    notes_style = ParagraphStyle(
        "Notes",
        parent=styles["BodyText"],
        fontSize=10,
        textColor="#64748b",
        leading=14,
        spaceBefore=18,
        leftIndent=12,
    )

    story: list = []
    for idx, slide in enumerate(presentation.slides):
        story.append(Paragraph(escape(slide.title or " "), title_style))
        body_html = escape(slide.body or "").replace("\n", "<br/>")
        story.append(Paragraph(body_html, body_style))
        if slide.speaker_notes:
            notes_html = "Notes: " + escape(slide.speaker_notes).replace("\n", "<br/>")
            story.append(Paragraph(notes_html, notes_style))
        story.append(Spacer(1, 12))
        if idx < len(presentation.slides) - 1:
            story.append(PageBreak())

    doc.build(story)
    return buffer.getvalue()


# ──────────────────────────────────────────────────────────────────────────────
# Dispatcher
# ──────────────────────────────────────────────────────────────────────────────

EXPORT_FORMATS = {
    "html": {
        "render": render_html,
        "content_type": "application/zip",
        "extension": "zip",
    },
    "pptx": {
        "render": render_pptx,
        "content_type": (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ),
        "extension": "pptx",
    },
    "pdf": {
        "render": render_pdf,
        "content_type": "application/pdf",
        "extension": "pdf",
    },
}
